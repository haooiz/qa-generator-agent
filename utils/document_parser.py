"""문서 파싱 유틸리티 (PDF, DOCX, PPTX, TXT)."""

from io import BytesIO

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PdfReader(BytesIO(file_bytes))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages).strip()


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = Document(BytesIO(file_bytes))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_text_from_txt(file_bytes: bytes) -> str:
    for encoding in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            return file_bytes.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace").strip()


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    slide_texts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if parts:
            slide_texts.append(f"[슬라이드 {slide_num}]\n" + "\n".join(parts))

    return "\n\n".join(slide_texts).strip()


def parse_document(filename: str, file_bytes: bytes) -> str:
    """파일명 확장자에 따라 텍스트를 추출합니다."""
    lower_name = filename.lower()

    if lower_name.endswith(".pdf"):
        text = extract_text_from_pdf(file_bytes)
    elif lower_name.endswith(".docx"):
        text = extract_text_from_docx(file_bytes)
    elif lower_name.endswith(".txt"):
        text = extract_text_from_txt(file_bytes)
    elif lower_name.endswith(".pptx"):
        text = extract_text_from_pptx(file_bytes)
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {filename}")

    if not text:
        raise ValueError("문서에서 텍스트를 추출할 수 없습니다.")

    return text
